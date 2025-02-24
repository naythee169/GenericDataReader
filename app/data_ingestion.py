import pandas as pd
import json
from pptx import Presentation
from PyPDF2 import PdfReader
import os

class DataIngestion:
    def __init__(self, data_dir=None):
        if data_dir is None:
            # Get the base directory: go up one level from the current file (inside app/)
            base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
            self.data_dir = os.path.join(base_dir, "data_files")
        else:
            self.data_dir = data_dir

    def read_csv(self, filename="sample.csv"):
        path = os.path.join(self.data_dir, filename)
        try:
            df = pd.read_csv(path)
            # Convert dataframe to list of dicts
            records = df.to_dict(orient="records")
            # Annotate each record with file type
            for rec in records:
                rec['file_type'] = 'csv'
            return records
        except Exception as e:
            print(f"Error reading CSV: {e}")
            return []

    def read_json(self, filename="sample.json"):
        path = os.path.join(self.data_dir, filename)
        try:
            with open(path, "r") as f:
                data = json.load(f)

            # Check if the data is a dict with a "companies" key.
            if isinstance(data, dict) and "companies" in data and isinstance(data["companies"], list):
                records = data["companies"]
                for rec in records:
                    if isinstance(rec, dict):
                        rec["file_type"] = "json"
                return records
            # If the data is already a list
            elif isinstance(data, list):
                for rec in data:
                    if isinstance(rec, dict):
                        rec["file_type"] = "json"
                return data
            else:
                # Fallback: wrap the data in a list.
                return [{"content": data, "file_type": "json"}]
        except Exception as e:
            print(f"Error reading JSON: {e}")
            return []


    def read_pptx(self, filename="sample.pptx"):
        path = os.path.join(self.data_dir, filename)
        try:
            prs = Presentation(path)
            records = []
            for slide in prs.slides:
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + " "
                records.append({"content": slide_text.strip(), "file_type": "pptx"})
            return records
        except Exception as e:
            print(f"Error reading PPTX: {e}")
            return []

    def read_pdf(self, filename="sample.pdf"):
        path = os.path.join(self.data_dir, filename)
        try:
            reader = PdfReader(path)
            content = ""
            for page in reader.pages:
                content += page.extract_text() + " "
            # For demo, return a single record with all text
            return [{"content": content.strip(), "file_type": "pdf"}]
        except Exception as e:
            print(f"Error reading PDF: {e}")
            return []

    def read_all(self):
        data = []
        data.extend(self.read_csv())
        data.extend(self.read_json())
        data.extend(self.read_pptx())
        data.extend(self.read_pdf())
        return data

    def read_by_type(self, file_type):
        if file_type == "csv":
            return self.read_csv()
        elif file_type == "json":
            return self.read_json()
        elif file_type == "pptx":
            return self.read_pptx()
        elif file_type == "pdf":
            return self.read_pdf()
        else:
            return []